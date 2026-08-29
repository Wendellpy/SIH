from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_coverage_slide():
    prs = Presentation()
    slide_layout = prs.slide_layouts[5] # Blank slide with title
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "PS Requirement Coverage: ULPIN Generation Cases"
    
    # 3-panel grid setup
    panel_width = Inches(3.0)
    panel_height = Inches(4.5)
    top = Inches(2.0)
    
    # Panel 1: Surface Parcel
    left1 = Inches(0.5)
    txBox1 = slide.shapes.add_textbox(left1, top, panel_width, panel_height)
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = "Surface Parcel (New)"
    p.font.bold = True
    p.font.size = Pt(20)
    
    p = tf1.add_paragraph()
    p.text = "ULPIN: MH13BOM04521873.S00-SURFACE"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    p = tf1.add_paragraph()
    p.text = "\nScale Achieved:\n8 real cadastral parcels processed across Mumbai Suburban & Mumbai City districts."
    p.font.size = Pt(12)
    
    p = tf1.add_paragraph()
    p.text = "\nStatus: Validated with Property Card generator. Slope-correction fallback added."
    p.font.size = Pt(12)
    
    # Panel 2: Multi-storey Building
    left2 = Inches(3.8)
    txBox2 = slide.shapes.add_textbox(left2, top, panel_width, panel_height)
    tf2 = txBox2.text_frame
    p = tf2.add_paragraph()
    p.text = "Multi-storey Building"
    p.font.bold = True
    p.font.size = Pt(20)
    
    p = tf2.add_paragraph()
    p.text = "ULPIN: MH13BOM04521873.A+02-201"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    p = tf2.add_paragraph()
    p.text = "\nCoverage:\nFully supported via Floorplan vectorization and 3D Extrusion pipeline."
    p.font.size = Pt(12)
    
    # Panel 3: Underground Infrastructure
    left3 = Inches(7.1)
    txBox3 = slide.shapes.add_textbox(left3, top, panel_width, panel_height)
    tf3 = txBox3.text_frame
    p = tf3.add_paragraph()
    p.text = "Underground Infra"
    p.font.bold = True
    p.font.size = Pt(20)
    
    p = tf3.add_paragraph()
    p.text = "ULPIN: MH13BOM04521873.U-01-WSUP12"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    p = tf3.add_paragraph()
    p.text = "\nCoverage:\nSupported via 3D Depth profiles and Utility clearance checking."
    p.font.size = Pt(12)
    
    prs.save("ps_requirement_coverage.pptx")
    print("Generated ps_requirement_coverage.pptx")

if __name__ == "__main__":
    create_coverage_slide()
